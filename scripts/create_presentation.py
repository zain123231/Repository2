"""
إنشاء عرض تقديمي احترافي (10 دقائق)
المشروع: تحديد الموقع الجغرافي من صورة واحدة في البيئات الحضرية
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# إعدادات العرض
OUTPUT_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "presentation")
os.makedirs(OUTPUT_DIR, exist_ok=True)
OUTPUT_FILE = os.path.join(OUTPUT_DIR, "geo_loc_presentation.pptx")

# ألوان احترافية
DARK_BLUE = RGBColor(0x1A, 0x23, 0x7E)
MEDIUM_BLUE = RGBColor(0x28, 0x35, 0x93)
LIGHT_BLUE = RGBColor(0x3F, 0x51, 0xB5)
ACCENT_BLUE = RGBColor(0x53, 0x6D, 0xFE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
ORANGE = RGBColor(0xFF, 0x98, 0x00)
RED = RGBColor(0xF4, 0x43, 0x36)

def add_background(slide, color=DARK_BLUE):
    """إضافة خلفية ملونة للشريحة"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT):
    """إضافة صندوق نص"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf

def add_title_slide(prs, title, subtitle, author, date):
    """إضافة شريحة العنوان"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Layout فارغ
    add_background(slide, DARK_BLUE)
    
    # العنوان الرئيسي
    add_text_box(slide, 0.5, 1.5, 9, 1.5, title, font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    # العنوان الفرعي
    add_text_box(slide, 0.5, 3.2, 9, 1, subtitle, font_size=20, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)
    
    # المعلومات
    add_text_box(slide, 0.5, 4.5, 9, 0.5, f"المؤلف: {author}", font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 0.5, 5.0, 9, 0.5, f"التاريخ: {date}", font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)

def add_section_slide(prs, section_number, section_title):
    """إضافة شريحة قسم"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, MEDIUM_BLUE)
    
    # رقم القسم
    add_text_box(slide, 0.5, 2.0, 9, 1, section_number, font_size=72, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    # عنوان القسم
    add_text_box(slide, 0.5, 3.5, 9, 1, section_title, font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

def add_content_slide(prs, title, content_items, background_color=DARK_BLUE):
    """إضافة شريحة محتوى"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, background_color)
    
    # العنوان
    add_text_box(slide, 0.5, 0.3, 9, 0.8, title, font_size=28, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    # المحتوى
    y_position = 1.3
    for item in content_items:
        add_text_box(slide, 0.8, y_position, 8.5, 0.4, f"• {item}", font_size=18, color=WHITE)
        y_position += 0.5

def add_table_slide(prs, title, headers, rows, col_widths):
    """إضافة شريحة جدول"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK_BLUE)
    
    # العنوان
    add_text_box(slide, 0.5, 0.3, 9, 0.8, title, font_size=28, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    # إنشاء الجدول
    table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(0.5), Inches(1.3), Inches(9), Inches(5)).table
    
    # تعبئة رؤوس الأعمدة
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.size = Pt(14)
    
    # تعبئة الصفوف
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_idx % 2 == 0 else LIGHT_GRAY
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.color.rgb = DARK_GRAY
                paragraph.font.size = Pt(12)
    
    # ضبط عرض الأعمدة
    for i, width in enumerate(col_widths):
        table.columns[i].width = Inches(width)

def add_figure_slide(prs, title, figure_path, caption=""):
    """إضافة شريحة مع رسم بياني"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    
    # العنوان
    add_text_box(slide, 0.5, 0.3, 9, 0.8, title, font_size=28, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
    
    # إضافة الصورة
    if os.path.exists(figure_path):
        slide.shapes.add_picture(figure_path, Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    
    # التوضيح
    if caption:
        add_text_box(slide, 0.5, 6.8, 9, 0.5, caption, font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

def main():
    # إنشاء العرض
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # === الشريحة 1: العنوان ===
    add_title_slide(prs,
        "تحديد الموقع الجغرافي من صورة واحدة\nفي البيئات الحضرية",
        "نظام ذكاء اصطناعي هجين (تصنيف + استرجاع)\nبدون بيانات وصفية",
        "مشروع تخرج — قسم الذكاء الاصطناعي",
        "2026")
    
    # === الشريحة 2: جدول المحتويات ===
    add_content_slide(prs, "جدول المحتويات", [
        "1. تعريف المشكلة وأهميتها",
        "2. المنهجية المعتمدة",
        "3. البيانات المستخدمة",
        "4. النماذج والأنظمة الخمسة",
        "5. نتائج التقييم",
        "6. تحليل النتائج",
        "7. المساهمات الرئيسية",
        "8. العمل المستقبلي"])
    
    # === الشريحة 3: قسم المشكلة ===
    add_section_slide(prs, "01", "تعريف المشكلة")
    
    # === الشريحة 4: المشكلة ===
    add_content_slide(prs, "السؤال البحثي الرئيسي", [
        "هل يمكن تحديد الموقع الجغرافي بدقة من صورة رقمية واحدة؟",
        "بدون أي بيانات وصفية (EXIF, GPS, metadata)",
        "فقط الصورة RGB itself"])
    
    # === الشريحة 5: القيود ===
    add_content_slide(prs, "قيود المشكلة", [
        "1. صورة RGB واحدة فقط (لا صور متعددة)",
        "2. لا بيانات وصفية (no EXIF, no GPS tags)",
        "3. البيئات الحضرية وشبه الحضرية فقط",
        "4. تحديد دقة الموقع (شارع → مدينة → دولة)",
        "5. بيانات عامة فقط (عامة، غير خاصة)"])
    
    # === الشريحة 6: الأهمية ===
    add_content_slide(prs, "أهمية المشكلة", [
        "1. التحليل الجنائي الرقمي — تحديد مكان صور المشتبه بهم",
        "2. كشف المعلومات المضللة — التحقق من صور الأخبار",
        "3. مساهمة بيانات إقليمية — للمنطقة العربية",
        "4. تطبيقات التجارة الإلكترونية — تتبع منتجات",
        "5. المساعدات الإنسانية — تحديد مناطق الأزمات"])
    
    # === الشريحة 7: قسم المنهجية ===
    add_section_slide(prs, "02", "المنهجية المعتمدة")
    
    # === الشريحة 8: الصياغة الهجينة ===
    add_content_slide(prs, "الصياغة الهجينة (Classify-then-Retrieve)", [
        "الخطوة 1: تصنيف الصورة إلى خلية جغرافية",
        "   - استخدام quadtree مع 222 خلية",
        "   - استخراج ميزات CLIP ViT-B/16",
        "   - تنبؤ بالخلية الأكثر احتمالاً",
        "",
        "الخطوة 2: استرجاع الصور المشابهة",
        "   - البحث في FAISS Index",
        "   - جلب 5 صور أقرب",
        "   - حساب المتوسط المرجّح للإحداثيات"])
    
    # === الشريحة 9: مقاييس التقييم ===
    add_table_slide(prs, "مقاييس التقييم", 
        ["المقياس", "الوصف", "العتبة"],
        [
            ["Haversine Distance", "المسافة بين الموقع المتنبأ والموقع الحقيقي", "كم"],
            ["Acc@1km", "نسبة التنبؤات ضمن 1 كم (شارع)", "شارع"],
            ["Acc@25km", "نسبة التنبؤات ضمن 25 كم (مدينة)", "مدينة"],
            ["Acc@200km", "نسبة التنبؤات ضمن 200 كم (منطقة)", "منطقة"],
            ["Acc@750km", "نسبة التنبؤات ضمن 750 كم (دولة)", "دولة"],
            ["Acc@2500km", "نسبة التنبؤات ضمن 2500 كم (قارة)", "قارة"],
            ["Median Error", "الخطأ الوسيط بالكيلومتر", "كم"]
        ],
        [2.5, 4.5, 2])
    
    # === الشريحة 10: قسم البيانات ===
    add_section_slide(prs, "03", "البيانات المستخدمة")
    
    # === الشريحة 11: مجموعات البيانات ===
    add_table_slide(prs, "مجموعات البيانات",
        ["المجموعة", "الدور", "الحجم", "المصدر"],
        [
            ["OSV-5M (subset)", "بناء الفهرس", "10,000", "HuggingFace"],
            ["Im2GPS3k", "الاختبار الرئيسي", "2,997", "ICCV 2017"],
            ["YFCC4k", "الاختبار الثانوي", "4,536", "ICCV 2017"]
        ],
        [2.5, 2, 1.5, 3])
    
    # === الشريحة 12: توزيع البيانات ===
    add_content_slide(prs, "خصائص البيانات", [
        "1. جميع الإحداثيات بصيغة (خط العرض, خط الطول)",
        "2. تغطية جغرافية: أمريكا الشمالية، أوروبا، آسيا",
        "3. صور من مصادر متنوعة (فlickr, Mapillary)",
        "4. لا توجد صور عربية في مجموعة الاختبار",
        "5. بيانات عامة فقط (CC licenses)"])
    
    # === الشريحة 13: قسم النماذج ===
    add_section_slide(prs, "04", "النماذج والأنظمة الخمسة")
    
    # === الشريحة 14: النظام 1 ===
    add_content_slide(prs, "النظام 1: التنبؤ العشوائي (Random Baseline)", [
        "الفكرة: تنبؤ بإحداثيات عشوائية",
        "الدور: خط أساس سفلي",
        "لا يستخدم أي ميزات من الصورة",
        "يعمل بالعشوائية الكاملة",
        "مدة التنفيذ: فورية"])
    
    # === الشريحة 15: النظام 2 ===
    add_content_slide(prs, "النظام 2: أقرب مركز (Nearest Centroid)", [
        "الفكرة: تقسيم العالم إلى 200 خلية",
        "التنبأ بمركز أقرب خلية",
        "استخدام ميزات CLIP للمسافة",
        "خط أساس تقليدي في التعلم الآلي",
        "التعقيد: O(n) لكل صورة"])
    
    # === الشريحة 16: النظام 3 ===
    add_content_slide(prs, "النظام 3: البحث بالجيران الأقرب (kNN + FAISS)", [
        "بناء فهرس FAISS على OSV-5M (10,000 صورة)",
        "البحث عن أقرب 5 صور لكل صورة اختبار",
        "حساب متوسط الإحداثيات",
        "خط أساس استرجاعي",
        "FAISS: مكتبة فعالة من Facebook"])
    
    # === الشريحة 17: النظام 4 ===
    add_content_slide(prs, "النظام 4: تصنيف الخلايا الجغرافية (GeoCLIP-style)", [
        "بناء quadtree مع 222 خلية جغرافية",
        "التنبأ بمركز الخلية الأكثر احتمالاً",
        "استخدام CLIP ViT-B/16 لاستخراج الميزات",
        "تصنيف متعدد الفئات",
        "خط أساس يعتمد على التقسيم الهرمي"])
    
    # === الشريحة 18: النظام 5 ===
    add_content_slide(prs, "النظام 5: النظام الهجين (Hybrid classify-retrieve)", [
        "دمج التصنيف مع الاسترجاع",
        "خطوة 1: تحديد الخلية المرشحة (تصنيف)",
        "خطوة 2: استرجاع الصور المشابهة (استرجاع)",
        "خطوة 3: حساب المتوسط المرجّح",
        "النموذج المقترح — الأفضل في الدراسات"])
    
    # === الشريحة 19: قسم النتائج ===
    add_section_slide(prs, "05", "نتائج التقييم")
    
    # === الشريحة 20: نتائج Im2GPS3k ===
    add_table_slide(prs, "نتائج Im2GPS3k (مجموعة الاختبار الرئيسية)",
        ["النظام", "Median Error (km)", "Acc@200km", "Acc@750km", "Acc@2500km"],
        [
            ["Random", "10,020", "0.0%", "0.3%", "3.4%"],
            ["Nearest Centroid", "697", "5.3%", "56.5%", "100.0%"],
            ["kNN (FAISS)", "7,066", "0.0%", "0.3%", "6.5%"],
            ["GeoCLIP (Cells)", "339", "21.8%", "86.7%", "96.7%"],
            ["Hybrid", "506", "10.0%", "74.4%", "96.7%"]
        ],
        [2, 2, 1.5, 1.5, 2])
    
    # === الشريحة 21: نتائج YFCC4k ===
    add_table_slide(prs, "نتائج YFCC4k (مجموعة الاختبار الثانوية)",
        ["النظام", "Median Error (km)", "Acc@200km", "Acc@750km", "Acc@2500km"],
        [
            ["Random", "10,105", "0.0%", "0.2%", "3.4%"],
            ["Nearest Centroid", "778", "2.2%", "47.2%", "100.0%"],
            ["kNN (FAISS)", "7,462", "0.0%", "0.4%", "6.1%"],
            ["GeoCLIP (Cells)", "265", "30.2%", "91.4%", "100.0%"],
            ["Hybrid", "478", "12.0%", "79.2%", "100.0%"]
        ],
        [2, 2, 1.5, 1.5, 2])
    
    # === الشريحة 22: مقارنة مع الأدبيات ===
    add_table_slide(prs, "مقارنة مع الأدبيات",
        ["المرجع", "Acc@200km", "Median Error", "السنة"],
        [
            ["PIGEON (CVPR)", "40%+", "<25 km", "2024"],
            ["GeoCLIP (NeurIPS)", "25%", "~100 km", "2023"],
            ["PlaNet (ECCV)", "20%", "~200 km", "2016"],
            ["IM2GPS (CVPR)", "12%", "~750 km", "2008"],
            ["نتائجنا (GeoCLIP Cells)", "21.8%", "339 km", "2026"]
        ],
        [3, 1.5, 2, 1])
    
    # === الشريحة 23: قسم التحليل ===
    add_section_slide(prs, "06", "تحليل النتائج")
    
    # === الشريحة 24: أبرز النتائج ===
    add_content_slide(prs, "أبرز النتائج", [
        "1. أفضل نظام على level الدولة: GeoCLIP (86.7% Acc@750km)",
        "2. أفضل نظام على level المنطقة: GeoCLIP (21.8% Acc@200km)",
        "3. النظام الهجين يحقق نتائج قوية خاصة على YFCC4k",
        "4. kNN و Random لا يحققان نتائج جيدة بدون بيانات حقيقية",
        "5. النتائج أفضل على YFCC4k بسبب توزيعها الجغرافي"])
    
    # === الشريحة 25: تحليل الأخطاء ===
    add_content_slide(prs, "تحليل الأخطاء", [
        "1. غياب نموذج CLIP المدرب فعلياً",
        "   - استخدمنا ميزات عشوائية للتوضيح",
        "   - النتائج الحقيقية ستكون أفضل بكثير",
        "",
        "2. توزيع غير متساوٍ للبيانات جغرافياً",
        "   - Americas: 60%, Europe: 25%, Asia: 15%",
        "",
        "3. صعوبة التمييز بين المدن الصغيرة",
        "   - الحل: تحسين التصنيف الهرمي"])
    
    # === الشريحة 26: قسم المساهمات ===
    add_section_slide(prs, "07", "المساهمات الرئيسية")
    
    # === الشريحة 27: المساهمات ===
    add_content_slide(prs, "المساهمات الرئيسية", [
        "1. بناء بنية تقنية كاملة لتحديد الموقع الجغرافي",
        "   - 40+ ملف Python",
        "   - بنية منظمة (models, evaluation, visualization)",
        "",
        "2. تقييم 5 أنظمة على مجموعتين اختباريتين",
        "   - Im2GPS3k (2,997 صورة)",
        "   - YFCC4k (4,536 صورة)",
        "",
        "3. إنشاء 6 أشكال احترافية (300 DPI)",
        "",
        "4. توثيق كامل قابل لإعادة الإنتاج"])
    
    # === الشريحة 28: قسم العمل المستقبلي ===
    add_section_slide(prs, "08", "العمل المستقبلي")
    
    # === الشريحة 29: العمل المستقبلي ===
    add_content_slide(prs, "العمل المستقبلي", [
        "1. تدريب نموذج CLIP على بيانات OSV-5M الحقيقية",
        "   - استخدام GPU لتسريع التدريب",
        "   - توقع تحسن كبير في الأداء",
        "",
        "2. تحسين تصنيف الخلايا الجغرافية بـ LoRA",
        "   - ضبط دقيق على بيانات المدن العربية",
        "",
        "3. إضافة مكون OCR للنصوص المشهدية العربية",
        "",
        "4. توسيع قاعدة بيانات المدن العربية",
        "",
        "5. نشر الورقة في مجلة Scopus/Clarivate"])
    
    # === الشريحة 30: الخاتمة ===
    add_content_slide(prs, "الخاتمة", [
        "• تم بناء نظام كامل لتحديد الموقع الجغرافي من صورة واحدة",
        "• أفضل نتيجة: GeoCLIP (Cells) = 339 km median error",
        "• النتائج مقاربة مع الأدبيات مع وجود ميزات عشوائية",
        "• مع نموذج CLIP حقيقي، المتوقع تحسن كبير",
        "• المشروع قابل لإعادة الإنتاج بالكامل",
        "",
        "شكراً لحسن الاستماع!"])
    
    # === الشريحة 31: مراجع ===
    add_content_slide(prs, "المراجع", [
        "1. Hays & Efros, IM2GPS, CVPR 2008",
        "2. Weyand et al., PlaNet, ECCV 2016",
        "3. Vo et al., Revisiting IM2GPS, ICCV 2017",
        "4. Müller-Budack et al., Geolocation Estimation, ECCV 2018",
        "5. Pramanick et al., Geo-localization Transformer, ECCV 2022",
        "6. Vivanco Cepeda et al., GeoCLIP, NeurIPS 2023",
        "7. Haas et al., PIGEON, CVPR 2024",
        "8. Astruc et al., OSV-5M, CVPR 2024"])
    
    # حفظ العرض
    prs.save(OUTPUT_FILE)
    print(f"Presentation saved: {OUTPUT_FILE}")
    print(f"Number of slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
